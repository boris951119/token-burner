# -*- coding: utf-8 -*-
r"""生成投资人汇报 PPT（v0.5.0-beta 版，聚焦：核心优势 / 竞品对比矩阵 / 商业化前景）。

运行：python scripts\make_investor_ppt.py
输出：investor_docs\token-burner_投资人汇报_v0.5.0-beta.pptx（独立文件夹存放，文件名含版本号）

数据口径（2026-09）：
- 产品事实：本仓库 v0.5.0-beta；pytest 826 passed / 7 skipped（共 833 项，0 失败）
- 竞品数据：Devin（thebestaitools/cognition.com）、MetaGPT（aiwiki/arXiv）、
  Claude Code/Cursor/Copilot（toolchase/aiunpacker/brightcoding，2026.05-08）
"""
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent.parent / "_docgen"))

from pptx import Presentation  # noqa: E402
from pptx.dml.color import RGBColor  # noqa: E402
from pptx.enum.shapes import MSO_SHAPE  # noqa: E402
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN  # noqa: E402
from pptx.util import Emu, Inches, Pt  # noqa: E402

ROOT = Path(__file__).resolve().parent.parent
OUT_DIR = ROOT / "investor_docs"
OUT_DIR.mkdir(exist_ok=True)
OUT = OUT_DIR / "token-burner_投资人汇报_v0.5.0-beta.pptx"

# ---- 主题色 ----
BG = RGBColor(0x10, 0x18, 0x28)
PANEL = RGBColor(0x1B, 0x25, 0x3A)
PANEL2 = RGBColor(0x24, 0x33, 0x52)      # 高亮面板（token-burner 列）
ACCENT = RGBColor(0xF5, 0x9E, 0x0B)
ACCENT2 = RGBColor(0xFF, 0x78, 0x47)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0xA8, 0xB3, 0xC5)
GREEN = RGBColor(0x34, 0xD3, 0x99)
RED = RGBColor(0xF8, 0x71, 0x71)
FONT = "微软雅黑"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def new_slide():
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    r.fill.solid()
    r.fill.fore_color.rgb = BG
    r.line.fill.background()
    r.shadow.inherit = False
    return s


def box(slide, l, t, w, h, fill=None, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.adjustments[0] = 0.06
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(1)
    sh.shadow.inherit = False
    return sh


def txt(slide, l, t, w, h, runs, size=18, color=WHITE, bold=False,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, leading=1.15):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    paras = runs if isinstance(runs, list) and runs and isinstance(runs[0], list) else [runs]
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = leading
        p.space_after = Pt(6)
        items = para if isinstance(para, list) else [(para, {})]
        for text, st in items:
            r = p.add_run()
            r.text = text
            f = r.font
            f.name = st.get("font", FONT)
            f.size = Pt(st.get("size", size))
            f.bold = st.get("bold", bold)
            f.color.rgb = st.get("color", color)
    return tb


def header(slide, kicker, title):
    box(slide, Inches(0.55), Inches(0.42), Inches(0.09), Inches(0.86), fill=ACCENT)
    txt(slide, Inches(0.85), Inches(0.38), Inches(11.9), Inches(0.34),
        kicker, size=13, color=ACCENT, bold=True)
    txt(slide, Inches(0.85), Inches(0.68), Inches(11.9), Inches(0.62),
        title, size=29, color=WHITE, bold=True)


def footer(slide, n):
    txt(slide, Inches(0.55), Inches(7.06), Inches(6), Inches(0.3),
        "token-burner · 文档消耗器 · v0.5.0-beta", size=10, color=GRAY)
    txt(slide, Inches(12.3), Inches(7.06), Inches(0.55), Inches(0.3),
        str(n), size=10, color=GRAY, align=PP_ALIGN.RIGHT)


def cell_run(cell, text, size, color, bold=False, align=PP_ALIGN.LEFT):
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = Inches(0.09)
    cell.margin_right = Inches(0.05)
    cell.margin_top = cell.margin_bottom = Inches(0.01)
    p = cell.text_frame.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color


# =========================================================================
# S1 封面
# =========================================================================
s = new_slide()
box(s, Inches(9.2), Inches(-1.6), Inches(7), Inches(7), fill=PANEL)
fire = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.6), Inches(1.6),
                          Inches(3.4), Inches(3.4))
fire.fill.solid(); fire.fill.fore_color.rgb = ACCENT2
fire.line.fill.background(); fire.shadow.inherit = False
txt(s, Inches(10.6), Inches(2.6), Inches(3.4), Inches(1.4),
    [[("燃", {"size": 52, "bold": True, "color": BG})]], align=PP_ALIGN.CENTER)
txt(s, Inches(0.9), Inches(1.95), Inches(8.6), Inches(0.4),
    "AI MULTI-AGENT PROJECT TEAM SYSTEM", size=14, color=ACCENT, bold=True)
txt(s, Inches(0.9), Inches(2.4), Inches(9.4), Inches(1.7),
    [[("token-burner", {"size": 54, "bold": True, "color": WHITE}),
      (" 文档消耗器", {"size": 40, "bold": True, "color": ACCENT})]])
txt(s, Inches(0.9), Inches(4.0), Inches(8.9), Inches(0.95),
    "成本可预算 · 质量可验证 · 过程可审计的 AI 多智能体软件工厂",
    size=19, color=GRAY)
box(s, Inches(0.9), Inches(5.15), Inches(2.9), Inches(0.55), fill=ACCENT)
txt(s, Inches(0.9), Inches(5.24), Inches(2.9), Inches(0.4),
    "投资人汇报 · 2026.09", size=15, color=BG, bold=True, align=PP_ALIGN.CENTER)
txt(s, Inches(0.9), Inches(6.45), Inches(9.2), Inches(0.4),
    "v0.5.0-beta 已交付  ｜ 833 项自动化测试 · 0 失败  ｜ Web / 桌面 / CLI / VS Code 四端形态",
    size=14, color=GRAY)

# =========================================================================
# S2 一页看懂
# =========================================================================
s = new_slide()
header(s, "01 · 电梯陈述", "一页看懂 token-burner")
box(s, Inches(0.55), Inches(1.72), Inches(12.23), Inches(1.1), fill=PANEL)
txt(s, Inches(0.95), Inches(1.95), Inches(11.5), Inches(0.7),
    [[("「不做又一个编码助手，做 AI 软件团队的", {"size": 21, "bold": True, "color": WHITE}),
      ("操作系统 + 财务总监", {"size": 21, "bold": True, "color": ACCENT}),
      ("。」", {"size": 21, "bold": True, "color": WHITE})]])
cards = [
    ("输入什么", "一句自然语言需求\n「帮我开发一个待办事项应用」"),
    ("输出什么", "一个可运行项目\n代码 + 测试 + 文档 + 成本报告"),
    ("怎么做到", "三模型 AI 团队协作\n六层确定性护栏全程约束"),
]
for i, (t1, t2) in enumerate(cards):
    x = Inches(0.55 + i * 4.18)
    box(s, x, Inches(3.05), Inches(3.9), Inches(1.7), fill=PANEL)
    txt(s, x + Inches(0.32), Inches(3.28), Inches(3.3), Inches(0.4),
        t1, size=16, bold=True, color=ACCENT)
    txt(s, x + Inches(0.32), Inches(3.75), Inches(3.35), Inches(0.9),
        [[(ln, {})] for ln in t2.split("\n")], size=14, color=WHITE, leading=1.3)
stats = [
    ("833", "项自动化测试（0 失败）"),
    ("4 端", "Web / 桌面 exe / CLI / VS Code"),
    ("200k", "token 任务预算上界（开工前锁定）"),
    ("3 模型", "强制互异 · 异构交叉验证"),
]
for i, (v, k) in enumerate(stats):
    x = Inches(0.55 + i * 3.12)
    box(s, x, Inches(5.05), Inches(2.9), Inches(1.5), fill=PANEL)
    txt(s, x, Inches(5.22), Inches(2.9), Inches(0.6), v, size=30, bold=True,
        color=GREEN, align=PP_ALIGN.CENTER)
    txt(s, x + Inches(0.15), Inches(5.9), Inches(2.6), Inches(0.55), k, size=12.5,
        color=GRAY, align=PP_ALIGN.CENTER, leading=1.2)
txt(s, Inches(0.55), Inches(6.72), Inches(12.2), Inches(0.4),
    [[("产品理念：", {"bold": True, "color": ACCENT, "size": 14.5}),
      ("决策归 LLM，校验与边界归程序——智能负责判断，规则负责把关。",
       {"color": GRAY, "size": 14.5})]])
footer(s, 2)

# =========================================================================
# S3 市场机会
# =========================================================================
s = new_slide()
header(s, "02 · 市场机会", "赛道已被验证，治理层无人占领")
vals = [
    ("GitHub Copilot", "470 万付费订阅", "约 2000 万用户，Fortune 100 覆盖 90%"),
    ("Cursor", "20 亿美元 ARR", "AI 原生 IDE，2026 年最热开发者工具"),
    ("Devin", "约 7300 万美元 ARR", "全自主 AI 工程师，已进入 Fortune 50"),
]
txt(s, Inches(0.55), Inches(1.62), Inches(12.2), Inches(0.4),
    "① 赛道验证：AI 编码是生成式 AI 商业化最快的赛道", size=15, bold=True, color=WHITE)
for i, (n, v, d) in enumerate(vals):
    x = Inches(0.55 + i * 4.18)
    box(s, x, Inches(2.1), Inches(3.9), Inches(1.45), fill=PANEL)
    txt(s, x + Inches(0.3), Inches(2.28), Inches(3.3), Inches(0.35),
        n, size=14, bold=True, color=GRAY)
    txt(s, x + Inches(0.3), Inches(2.62), Inches(3.4), Inches(0.45),
        v, size=21, bold=True, color=GREEN)
    txt(s, x + Inches(0.3), Inches(3.1), Inches(3.4), Inches(0.35),
        d, size=11.5, color=GRAY)
txt(s, Inches(0.55), Inches(3.85), Inches(12.2), Inches(0.4),
    "② 范式转移：从「单点补全」到「团队交付」——但企业采用有三大阻力",
    size=15, bold=True, color=WHITE)
gaps = [
    ("成本不可预估", "ACU/额度事后计费，账单无法预算管理"),
    ("数据必须出域", "代码与过程交给云端 SaaS，合规受阻"),
    ("模型被锁定", "绑定单一厂商，无法按任务组合最优模型"),
]
for i, (t1, t2) in enumerate(gaps):
    x = Inches(0.55 + i * 4.18)
    box(s, x, Inches(4.32), Inches(3.9), Inches(1.35), fill=PANEL)
    box(s, x, Inches(4.32), Inches(3.9), Inches(0.1), fill=ACCENT2)
    txt(s, x + Inches(0.3), Inches(4.55), Inches(3.3), Inches(0.4),
        t1, size=16, bold=True, color=ACCENT2)
    txt(s, x + Inches(0.3), Inches(5.0), Inches(3.35), Inches(0.6),
        t2, size=12.5, color=GRAY, leading=1.25)
box(s, Inches(0.55), Inches(5.95), Inches(12.23), Inches(0.85), fill=PANEL2)
txt(s, Inches(0.95), Inches(6.13), Inches(11.5), Inches(0.55),
    [[("空白地带：", {"bold": True, "color": ACCENT, "size": 16}),
      ("企业关注点已从「能不能用」转向「预算可控 / 质量可验证 / 过程可审计」——"
       "多智能体编排 + 成本治理层是尚未被占领的下一层价值。",
       {"color": WHITE, "size": 16})]])
footer(s, 3)

# =========================================================================
# S4 解决方案：工作管线
# =========================================================================
s = new_slide()
header(s, "03 · 解决方案", "一支被工程纪律约束的 AI 项目团队")
stages = [
    ("需求", "不可信输入\n边界包裹"),
    ("评估路由", "三分类\n难度分 1-10"),
    ("组队", "三模型互异\nPM+Dev+Test"),
    ("方案讨论", "双评审打分\n≤3 轮收敛"),
    ("spec 确认", "用户确认\n≤3 次合并"),
    ("模块化生成", "接口契约\nAST 门禁"),
    ("双模验证", "静态+沙箱执行\n修复 ≤5 轮"),
    ("交付", "代码+测试+文档\n+成本报告"),
]
for i, (t1, t2) in enumerate(stages):
    x = 0.55 + i * 1.54
    box(s, Inches(x), Inches(1.95), Inches(1.38), Inches(1.7), fill=PANEL)
    box(s, Inches(x), Inches(1.95), Inches(1.38), Inches(0.08), fill=ACCENT)
    txt(s, Inches(x + 0.06), Inches(2.16), Inches(1.26), Inches(0.4),
        t1, size=13.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, Inches(x + 0.05), Inches(2.6), Inches(1.28), Inches(0.95),
        [[(ln, {})] for ln in t2.split("\n")], size=10, color=GRAY,
        align=PP_ALIGN.CENTER, leading=1.25)
    if i < len(stages) - 1:
        txt(s, Inches(x + 1.36), Inches(2.5), Inches(0.2), Inches(0.4),
            "›", size=18, bold=True, color=ACCENT)
hls = [
    ("六层护栏全程约束", "讨论不发散、修复不死循环、输出不倾泻——"
     "每一步都在预算闸门内"),
    ("三模型异构互验", "架构 / 实现 / 测试三视角独立校验，"
     "从结构上杜绝同源盲区"),
    ("中断可恢复 · 全程落盘", "进程崩溃重启即续跑；逐调用 token、"
     "讨论过程、修复历史全量审计"),
]
for i, (t1, t2) in enumerate(hls):
    y = Inches(4.05 + i * 0.78)
    box(s, Inches(0.55), y, Inches(12.23), Inches(0.68), fill=PANEL)
    box(s, Inches(0.55), y, Inches(0.1), Inches(0.68), fill=GREEN)
    txt(s, Inches(0.9), y + Inches(0.07), Inches(3.3), Inches(0.5),
        t1, size=14.5, bold=True, color=WHITE)
    txt(s, Inches(4.35), y + Inches(0.1), Inches(8.2), Inches(0.5),
        t2, size=13, color=GRAY)
txt(s, Inches(0.55), Inches(6.55), Inches(12.2), Inches(0.4),
    [[("模型无关接入：", {"bold": True, "color": ACCENT, "size": 14}),
      ("基于 litellm——OpenAI / Anthropic / DeepSeek / Gemini / 智谱 GLM 任意组合，自带密钥（BYOK）。",
       {"color": GRAY, "size": 14})]])
footer(s, 4)

# =========================================================================
# S5 核心优势① 六层成本护栏
# =========================================================================
s = new_slide()
header(s, "04 · 核心优势 ①", "六层确定性成本护栏：把「无限烧钱」变成「预算工程」")
rows = [
    ("层", "护栏", "默认参数", "达标行为"),
    ("0", "单任务总预算（总闸）", "200k tokens · 自动模式 ×2.5",
     "≥90% 省 token 模式；超预算中止并落盘止损清单"),
    ("1", "讨论轮数上限", "3 轮", "主 LLM 直接产出收敛 spec"),
    ("2", "单轮输出上限", "8k tokens", "截断按未完成处理，分块续写（≤2 次）"),
    ("3", "讨论循环检测", "Jaccard 0.9 / embedding 0.85 / 3 次",
     "冻结副 LLM 发言权，主 LLM 收权裁决"),
    ("4", "修复循环上限", "5 轮", "输出「已知问题与降级方案」"),
    ("5", "spec 确认收敛", "3 次", "合并意见输出最终 spec"),
]
tbl = s.shapes.add_table(7, 4, Inches(0.55), Inches(1.82),
                         Inches(12.23), Inches(3.9)).table
tbl.columns[0].width = Inches(0.6)
tbl.columns[1].width = Inches(3.1)
tbl.columns[2].width = Inches(3.9)
tbl.columns[3].width = Inches(4.63)
for ri, row in enumerate(rows):
    for ci, val in enumerate(row):
        cell = tbl.cell(ri, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT if ri == 0 else (
            PANEL if ri % 2 else RGBColor(0x16, 0x20, 0x33))
        cell_run(cell, val, 13 if ri else 14,
                 BG if ri == 0 else (ACCENT if ci == 0 else WHITE),
                 bold=ri == 0 or ci == 0)
box(s, Inches(0.55), Inches(5.95), Inches(12.23), Inches(0.85), fill=PANEL2)
txt(s, Inches(0.95), Inches(6.1), Inches(11.5), Inches(0.6),
    [[("代际差异：", {"bold": True, "color": ACCENT, "size": 15}),
      ("竞品的成本控制是「计费提醒」（事后账单），我们的是「预算闸门」（开工前锁定上界、"
       "超支即熔断）——这是采购与财务岗位听得懂的采购语言。",
       {"color": WHITE, "size": 15})]])
footer(s, 5)

# =========================================================================
# S6 核心优势② 质量是机制
# =========================================================================
s = new_slide()
header(s, "05 · 核心优势 ②", "质量不是提示词，是机制")
cards = [
    ("三模型强制互异", [
        "主 LLM（架构师）+ 开发副 + 测试副，模型必须互不相同",
        "异构视角交叉验证——单模型产品结构上无法复制",
        "双评审打分制：评审 ≥8 分建议自动采纳"]),
    ("AST 接口门禁", [
        "接口地图 = 模块间唯一合法基线（imports/exports/public_api）",
        "引用抽取差异判定，违规即阻断进入修复循环",
        "零执行、零 LLM 调用——唯一不耗 token 的纪律执行器"]),
    ("循环检测双保险", [
        "第一道：Jaccard 文本相似度（零成本拦截字面重复）",
        "第二道：embedding 语义相似度（增量缓存控成本）",
        "重复 ≥3 次自动冻结发言权，主 LLM 收权裁决"]),
    ("双模执行验证", [
        "安全审阅模式（默认）：静态检查 + LLM 审查，不执行任何代码",
        "自动验证模式：Docker 沙箱真实执行（512MB/CPU1/只读/无网络/30s 熔断）",
        "Docker 缺失自动降级进程模式，危险 API 预扫描始终生效"]),
]
for i, (t1, items) in enumerate(cards):
    x = Inches(0.55 + (i % 2) * 6.27)
    y = Inches(1.9 + (i // 2) * 2.42)
    box(s, x, y, Inches(5.95), Inches(2.2), fill=PANEL)
    box(s, x, y, Inches(0.1), Inches(2.2), fill=GREEN if i % 2 else ACCENT)
    txt(s, x + Inches(0.32), y + Inches(0.16), Inches(5.4), Inches(0.42),
        t1, size=17, bold=True, color=WHITE)
    paras = [[("▪ ", {"color": ACCENT2, "size": 12.5}),
              (it, {"color": GRAY, "size": 12.5})] for it in items]
    txt(s, x + Inches(0.32), y + Inches(0.62), Inches(5.35), Inches(1.5),
        paras, leading=1.28)
footer(s, 6)

# =========================================================================
# S7 核心优势③ 企业级基因
# =========================================================================
s = new_slide()
header(s, "06 · 核心优势 ③", "为企业采用而生的架构基因")
feats = [
    ("私有化部署", "本地运行，需求文本、生成代码、讨论过程全程不出域——"
     "对数据敏感行业（金融/政务/军工配套）是硬性门槛，云端 SaaS 结构上无法满足。"),
    ("模型中立 BYOK", "litellm 统一接入，密钥自带、按需切换——"
     "规避供应商锁定与涨价风险；模型市场越碎片化，中立编排层价值越大。"),
    ("全链路审计", "逐调用 token 落盘、按模型/阶段/档位成本分解、Agent 对话流图、"
     "修复历史、中断快照——合规与审计部门可直接取证。"),
    ("工程成熟度", "833 项测试回归网、异步任务 + SSE 实时进度、任务取消与僵尸清理、"
     "中断恢复、GitHub Actions 自动发布流水线——不是 demo，是可交付软件。"),
]
for i, (t1, t2) in enumerate(feats):
    x = Inches(0.55 + (i % 2) * 6.27)
    y = Inches(1.95 + (i // 2) * 2.15)
    box(s, x, y, Inches(5.95), Inches(1.95), fill=PANEL)
    box(s, x + Inches(0.3), y + Inches(0.28), Inches(0.55), Inches(0.55), fill=ACCENT)
    txt(s, x + Inches(0.3), y + Inches(0.36), Inches(0.55), Inches(0.45),
        str(i + 1), size=20, bold=True, color=BG, align=PP_ALIGN.CENTER)
    txt(s, x + Inches(1.05), y + Inches(0.3), Inches(4.6), Inches(0.5),
        t1, size=17.5, bold=True, color=WHITE)
    txt(s, x + Inches(1.05), y + Inches(0.85), Inches(4.7), Inches(1.0),
        t2, size=12.5, color=GRAY, leading=1.3)
box(s, Inches(0.55), Inches(6.35), Inches(12.23), Inches(0.55), fill=PANEL2)
txt(s, Inches(0.95), Inches(6.46), Inches(11.5), Inches(0.4),
    [[("四端形态降低采用门槛：", {"bold": True, "color": ACCENT, "size": 14}),
      ("Web 工作台 / 双击即用桌面 exe / CLI / VS Code 插件——同一内核，覆盖从小白到开发者。",
       {"color": WHITE, "size": 14})]])
footer(s, 7)

# =========================================================================
# S8 竞品对比矩阵（核心页）
# =========================================================================
s = new_slide()
header(s, "07 · 竞品对比矩阵", "唯一同时满足四项企业级约束的产品")
matrix = [
    ("对比维度", "token-burner", "Devin", "MetaGPT / ChatDev", "Copilot / Cursor / CC"),
    ("多智能体架构", "三模型强制互异互验", "单代理 + 模型路由", "多角色协作", "单模型代理"),
    ("成本治理", "六层护栏 + 预算总闸", "ACU 事后计费", "无", "订阅额度制"),
    ("开工前成本可预估", "✓ token 上界锁定", "✗ 事后账单", "✗", "✗"),
    ("验证闭环", "静态 + AST 门禁 + 双模执行", "沙箱自测 + PR", "生成测试文件", "开发者自测"),
    ("执行沙箱", "Docker 配额沙箱", "云端 VM", "无", "本地 / 受限"),
    ("模型中立 BYOK", "✓ litellm 任意组合", "✗ 厂商锁定", "✓", "✗"),
    ("私有化部署", "✓ 数据不出域", "企业版 VPC", "✓ 需自建", "✗"),
    ("过程可审计", "逐调用 + 对话流 + 快照", "平台内日志", "无", "无"),
    ("使用形态", "四端（含双击 exe）", "云端 SaaS + IDE", "开发者自建框架", "IDE / 终端"),
    ("单任务成本", "按量 ≈ $0.2–2（BYOK）", "$11–45（ACU）", "自付 API", "订阅内"),
]
tbl = s.shapes.add_table(11, 5, Inches(0.55), Inches(1.78),
                         Inches(12.23), Inches(4.85)).table
tbl.columns[0].width = Inches(1.85)
tbl.columns[1].width = Inches(2.75)
tbl.columns[2].width = Inches(2.45)
tbl.columns[3].width = Inches(2.45)
tbl.columns[4].width = Inches(2.73)
for ri, row in enumerate(matrix):
    for ci, val in enumerate(row):
        cell = tbl.cell(ri, ci)
        cell.fill.solid()
        if ri == 0:
            cell.fill.fore_color.rgb = ACCENT if ci == 1 else PANEL
        elif ci == 1:
            cell.fill.fore_color.rgb = PANEL2
        else:
            cell.fill.fore_color.rgb = (
                RGBColor(0x16, 0x20, 0x33) if ri % 2 else PANEL)
        if ri == 0:
            color = BG if ci == 1 else WHITE
            cell_run(cell, val, 10.5, color, bold=True)
        else:
            color = WHITE if ci <= 1 else GRAY
            if val.startswith("✓"):
                color = GREEN
            elif val.startswith("✗"):
                color = RED
            cell_run(cell, val, 9.5, color, bold=ci == 1)
txt(s, Inches(0.55), Inches(6.75), Inches(12.2), Inches(0.4),
    [[("注：竞品信息时点 2026.05–08（来源见附录）；", {"color": GRAY, "size": 11}),
      ("四项 ✓（可预估 / BYOK / 私有化 / 审计）在同一产品上同时成立者，市场上仅此一家。",
       {"color": ACCENT, "size": 11.5, "bold": True})]])
footer(s, 8)

# =========================================================================
# S9 竞品逐项分析
# =========================================================================
s = new_slide()
header(s, "08 · 竞品逐项分析", "他们强在哪，空在哪")
comps = [
    ("Devin（Cognition）", "$20–500/月 + ACU 计量",
     ["全自主端到端交付 PR", "Fortune 50 渗透 / FedRAMP"],
     ["单任务实际 $11–45，事前不可预算", "复杂多文件任务成功率仅 14–15%",
      "闭源 SaaS：代码与过程出域", "模型路由黑盒，厂商锁定"]),
    ("MetaGPT / ChatDev", "开源框架（69k star）",
     ["「Code = SOP(Team)」理念成熟", "MIT / Apache 开源，社区活跃"],
     ["无任何成本治理与预算闸门", "无产品化：无桌面 / 插件 / 恢复 / 审计",
      "质量约束靠提示词，非机制"]),
    ("Copilot / Cursor / Claude Code", "$10–200/月",
     ["编辑体验与模型能力最强（CC SWE-bench 80.8%）", "Copilot 470 万付费订阅验证付费意愿"],
     ["服务「写代码的人」，非项目交付", "单模型：无跨模型互验",
      "无任务级预算治理与审计"]),
]
for i, (name, price, pros, cons) in enumerate(comps):
    x = Inches(0.55 + i * 4.18)
    box(s, x, Inches(1.85), Inches(3.9), Inches(4.35), fill=PANEL)
    txt(s, x + Inches(0.28), Inches(2.05), Inches(3.35), Inches(0.42),
        name, size=15.5, bold=True, color=WHITE)
    txt(s, x + Inches(0.28), Inches(2.48), Inches(3.35), Inches(0.32),
        price, size=11.5, color=ACCENT, bold=True)
    paras = [[("✓ ", {"color": GREEN, "size": 11.5}),
              (p, {"color": GRAY, "size": 11.5})] for p in pros]
    txt(s, x + Inches(0.28), Inches(2.92), Inches(3.35), Inches(1.0),
        paras, leading=1.25)
    paras = [[("✗ ", {"color": RED, "size": 11.5}),
              (c, {"color": WHITE, "size": 11.5})] for c in cons]
    txt(s, x + Inches(0.28), Inches(4.05), Inches(3.35), Inches(2.0),
        paras, leading=1.3)
box(s, Inches(0.55), Inches(6.4), Inches(12.23), Inches(0.55), fill=PANEL2)
txt(s, Inches(0.95), Inches(6.51), Inches(11.5), Inches(0.4),
    [[("共同空白 = 我们的生态位：", {"bold": True, "color": ACCENT, "size": 14.5}),
      ("不与 Copilot 拼编辑体验、不与 Devin 拼企业生态——用「预算工程 + 私有化」占领无人区。",
       {"color": WHITE, "size": 14.5})]])
footer(s, 9)

# =========================================================================
# S10 商业化① 定价模型与成本优势
# =========================================================================
s = new_slide()
header(s, "09 · 商业化前景 ①", "三条收入线 + 一个数量级的成本优势")
lines = [
    ("SaaS 订阅", "按任务预算档位订阅——护栏 / 看板 / 审计为付费核心",
     "面向中小团队，最快起量"),
    ("私有化部署", "按席位 / 年授权，数据全程不出域",
     "面向合规敏感企业，客单价最高"),
    ("API 计量", "多智能体内核作为 PnP 组件按调用量计费",
     "面向集成商与开发者生态"),
]
for i, (t1, t2, t3) in enumerate(lines):
    y = Inches(1.9 + i * 1.18)
    box(s, Inches(0.55), y, Inches(6.6), Inches(1.02), fill=PANEL)
    box(s, Inches(0.55), y, Inches(0.1), Inches(1.02), fill=ACCENT)
    txt(s, Inches(0.9), y + Inches(0.12), Inches(2.2), Inches(0.4),
        t1, size=15.5, bold=True, color=WHITE)
    txt(s, Inches(0.9), y + Inches(0.54), Inches(6.0), Inches(0.4),
        t2, size=12, color=GRAY)
    txt(s, Inches(3.2), y + Inches(0.13), Inches(3.8), Inches(0.35),
        t3, size=11.5, color=ACCENT, bold=True, align=PP_ALIGN.RIGHT)
box(s, Inches(7.45), Inches(1.9), Inches(5.33), Inches(3.5), fill=PANEL)
txt(s, Inches(7.75), Inches(2.1), Inches(4.8), Inches(0.4),
    "同级任务成本对比（估算）", size=15, bold=True, color=WHITE)
txt(s, Inches(7.75), Inches(2.62), Inches(4.8), Inches(0.9),
    [[("token-burner", {"size": 13, "color": GRAY})],
     [("$0.2 – $2", {"size": 34, "bold": True, "color": GREEN})],
     [("200k 预算按主流模型公开价 · BYOK · 上界开工前锁定",
       {"size": 10.5, "color": GRAY})]], leading=1.15)
txt(s, Inches(7.75), Inches(4.1), Inches(4.8), Inches(0.9),
    [[("Devin（ACU 计量）", {"size": 13, "color": GRAY})],
     [("$11 – $45", {"size": 34, "bold": True, "color": RED})],
     [("中等复杂任务 5–20 ACU · 另加 $20–500 月费",
       {"size": 10.5, "color": GRAY})]], leading=1.15)
txt(s, Inches(7.75), Inches(5.02), Inches(4.8), Inches(0.35),
    "≈ 5–20 倍单任务成本差", size=15, bold=True, color=ACCENT,
    align=PP_ALIGN.CENTER)
box(s, Inches(0.55), Inches(5.62), Inches(12.23), Inches(1.1), fill=PANEL2)
txt(s, Inches(0.95), Inches(5.8), Inches(11.5), Inches(0.8),
    [[("毛利逻辑：", {"bold": True, "color": ACCENT, "size": 14.5}),
      ("治理能力（护栏 / 看板 / 审计 / 沙箱）是付费点而非 token 转售——", {"color": WHITE, "size": 14.5})],
     [("BYOK 模式下 token 成本透明转嫁（零垫资风险），托管模式可打包加价（参考行业 20–30% 毛利）。",
       {"color": GRAY, "size": 13.5})]], leading=1.3)
footer(s, 10)

# =========================================================================
# S11 商业化② 客群与 GTM
# =========================================================================
s = new_slide()
header(s, "10 · 商业化前景 ②", "目标客群与市场进入路径")
segs = [
    ("中小软件团队 / 外包公司", "降本刚需 + 账单必须可预算；"
     "一句需求出项目的交付模式直接替代低码外包"),
    ("企业 IT（合规敏感行业）", "数据不出域 + 全链路审计是硬门槛；"
     "私有化部署高客单价"),
    ("AI 集成商 / 开发者", "API / CLI 嵌入自有工作流，"
     "内核作为 PnP 组件分发"),
]
for i, (t1, t2) in enumerate(segs):
    y = Inches(1.9 + i * 1.32)
    box(s, Inches(0.55), y, Inches(5.6), Inches(1.16), fill=PANEL)
    box(s, Inches(0.55), y, Inches(0.1), Inches(1.16), fill=GREEN)
    txt(s, Inches(0.9), y + Inches(0.13), Inches(5.0), Inches(0.4),
        t1, size=15, bold=True, color=WHITE)
    txt(s, Inches(0.9), y + Inches(0.56), Inches(5.05), Inches(0.55),
        t2, size=12, color=GRAY, leading=1.25)
gtm = [
    ("楔子切入", "以「成本治理」为差异化卖点获取预算敏感种子客群"),
    ("数据沉淀", "真实任务成功率 / 成本节省数据持续积累（内置 A/B 框架）"),
    ("信誉建立", "公开基准评测（SWE-bench 子集）+ 代码签名 + 案例背书"),
    ("客单放大", "私有化部署 + 团队席位 + 生态 API 三级放大"),
]
txt(s, Inches(6.45), Inches(1.62), Inches(6.3), Inches(0.4),
    "GTM 四步走", size=15, bold=True, color=ACCENT)
for i, (t1, t2) in enumerate(gtm):
    y = Inches(2.08 + i * 1.02)
    box(s, Inches(6.45), y, Inches(6.33), Inches(0.88), fill=PANEL)
    box(s, Inches(6.45), y, Inches(0.42), Inches(0.88), fill=ACCENT)
    txt(s, Inches(6.45), y + Inches(0.24), Inches(0.42), Inches(0.4),
        str(i + 1), size=16, bold=True, color=BG, align=PP_ALIGN.CENTER)
    txt(s, Inches(7.05), y + Inches(0.1), Inches(5.5), Inches(0.35),
        t1, size=13.5, bold=True, color=WHITE)
    txt(s, Inches(7.05), y + Inches(0.46), Inches(5.6), Inches(0.35),
        t2, size=11.5, color=GRAY)
box(s, Inches(0.55), Inches(6.28), Inches(12.23), Inches(0.6), fill=PANEL2)
txt(s, Inches(0.95), Inches(6.4), Inches(11.5), Inches(0.4),
    [[("北极星指标：", {"bold": True, "color": ACCENT, "size": 14}),
      ("真实任务成功率 · 单任务成本节省比例 · 预算命中率（超支率 <10%）",
       {"color": WHITE, "size": 14})]])
footer(s, 11)

# =========================================================================
# S12 路线图
# =========================================================================
s = new_slide()
header(s, "11 · 路线图与里程碑", "MVP 已交付，商业化窗口清晰")
phases = [
    ("当前 · v0.5.0-beta", "2026 Q3", GREEN,
     ["833 项测试 / 四端形态 / Docker 沙箱", "Researcher 调研 / 异步任务 / 发布流水线",
      "已具备：真实任务跑通全链路"]),
    ("Beta 强化期", "2026 Q4", ACCENT,
     ["真实样本积累与成功率统计", "SWE-bench 子集公开评测",
      "代码签名 + 种子客户试点（付费验证）"]),
    ("v1.0 商业化", "2027 H1", ACCENT2,
     ["本地知识库检索（RAG）", "副 LLM 能力自动匹配",
      "多人协作与权限管理", "SaaS 订阅 + 私有化双轨收费"]),
    ("生态扩展", "2027 H2+", WHITE,
     ["外部工具 / MCP 开放", "Agent 市场与模板生态",
      "企业级治理（RBAC / 审计合规包）"]),
]
for i, (t1, t2, c, items) in enumerate(phases):
    x = Inches(0.55 + i * 3.12)
    box(s, x, Inches(1.95), Inches(2.9), Inches(4.15), fill=PANEL)
    box(s, x, Inches(1.95), Inches(2.9), Inches(0.14), fill=c)
    txt(s, x + Inches(0.25), Inches(2.28), Inches(2.45), Inches(0.72),
        t1, size=15.5, bold=True, color=c, leading=1.15)
    txt(s, x + Inches(0.25), Inches(3.0), Inches(2.4), Inches(0.35),
        t2, size=12, color=GRAY)
    paras = [[("▪ ", {"color": c, "size": 11.5}),
              (it, {"color": GRAY, "size": 11.5})] for it in items]
    txt(s, x + Inches(0.25), Inches(3.45), Inches(2.45), Inches(2.5),
        paras, leading=1.35)
txt(s, Inches(0.55), Inches(6.35), Inches(12.2), Inches(0.5),
    [[("里程碑意义：", {"bold": True, "color": ACCENT, "size": 14}),
      ("技术风险已在前两版消化完毕（每版全量回归）；Beta 期核心任务是商业化验证而非补功能。",
       {"color": GRAY, "size": 14})]])
footer(s, 12)

# =========================================================================
# S13 风险与对策
# =========================================================================
s = new_slide()
header(s, "12 · 风险与对策", "诚实面对差距，每条都有解")
risks = [
    ("品牌与市场认知缺失",
     "Devin 已入 Fortune 50 与联邦市场；本项目无认知度",
     "以「成本治理 + 私有化」楔子切入预算敏感客群，避开正面竞争"),
    ("尚无公开基准成绩",
     "竞品有 SWE-bench 等第三方数据（CC 80.8%）",
     "Beta 期在 SWE-bench 子集跑分并公开方法论与成本对照"),
    ("真实成功率待大样本验证",
     "833 项测试是 stub 回归网，非真实成功率承诺",
     "内置 A/B 框架 + 性能基线脚本持续积累；如实呈现降级方案"),
    ("基座模型能力快速迭代",
     "模型变强可能稀释编排层价值",
     "壁垒锚定确定性护栏与审计合规——模型越强，企业对治理需求越大"),
    ("分发与信任摩擦",
     "未签名 exe 触发部分安全软件误报",
     "发布流水线已就绪，商业化前采购代码签名证书"),
]
tbl = s.shapes.add_table(6, 3, Inches(0.55), Inches(1.85),
                         Inches(12.23), Inches(4.7)).table
tbl.columns[0].width = Inches(2.6)
tbl.columns[1].width = Inches(4.1)
tbl.columns[2].width = Inches(5.53)
head = ("风险", "客观描述", "对策")
for ci, val in enumerate(head):
    cell = tbl.cell(0, ci)
    cell.fill.solid()
    cell.fill.fore_color.rgb = ACCENT
    cell_run(cell, val, 14, BG, bold=True)
for ri, (r1, r2, r3) in enumerate(risks, start=1):
    for ci, val in enumerate((r1, r2, r3)):
        cell = tbl.cell(ri, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = PANEL if ri % 2 else RGBColor(0x16, 0x20, 0x33)
        color = ACCENT2 if ci == 0 else (GRAY if ci == 1 else WHITE)
        cell_run(cell, val, 11.5, color, bold=ci == 0)
txt(s, Inches(0.55), Inches(6.72), Inches(12.2), Inches(0.4),
    [[("结构性差异难以复制、市场性短板可以系统性弥补——", {"color": GRAY, "size": 13}),
      ("这是当前阶段的风险收益比。", {"color": ACCENT, "size": 13, "bold": True})]])
footer(s, 13)

# =========================================================================
# S14 结尾
# =========================================================================
s = new_slide()
box(s, Inches(0), Inches(1.9), SW, Inches(3.2), fill=PANEL)
txt(s, Inches(0), Inches(2.3), SW, Inches(0.9),
    "Set your tokens on fire — 有价值地燃烧", size=34, bold=True,
    color=WHITE, align=PP_ALIGN.CENTER)
txt(s, Inches(0), Inches(3.35), SW, Inches(0.6),
    "成本可预算 · 质量可验证 · 过程可审计的 AI 多智能体软件工厂",
    size=18, color=ACCENT, align=PP_ALIGN.CENTER)
topics = ["融资规划与资金用途", "种子客户与试点合作", "公开基准评测计划"]
for i, t in enumerate(topics):
    x = Inches(2.05 + i * 3.3)
    box(s, x, Inches(5.35), Inches(3.0), Inches(0.62), fill=PANEL2)
    txt(s, x, Inches(5.5), Inches(3.0), Inches(0.4),
        t, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, Inches(0), Inches(6.35), SW, Inches(0.5),
    "谢谢 · 期待与您深入交流", size=17, color=GRAY, align=PP_ALIGN.CENTER)

prs.save(OUT)
print("saved:", OUT, "| slides:", len(prs.slides.__iter__.__self__._sldIdLst))
