# -*- coding: utf-8 -*-
r"""生成投资人材料：PPT（pptx）+ 可行性研究报告（docx）。

运行：python scripts\make_investor_docs.py
输出：investor_docs\token-burner_投资人汇报.pptx
      investor_docs\token-burner_可行性研究报告.docx
"""
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent.parent / "_docgen"))

from docx import Document  # noqa: E402
from docx.enum.text import WD_ALIGN_PARAGRAPH  # noqa: E402
from docx.oxml.ns import qn  # noqa: E402
from pptx import Presentation  # noqa: E402
from pptx.dml.color import RGBColor  # noqa: E402
from pptx.enum.shapes import MSO_SHAPE  # noqa: E402
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN  # noqa: E402
from pptx.util import Emu, Inches, Pt  # noqa: E402

ROOT = Path(__file__).resolve().parent.parent
OUT = ROOT / "investor_docs"
OUT.mkdir(exist_ok=True)

# ---- 主题色（深空蓝 + 火焰橙）----
BG = RGBColor(0x10, 0x18, 0x28)
PANEL = RGBColor(0x1B, 0x25, 0x3A)
ACCENT = RGBColor(0xF5, 0x9E, 0x0B)
ACCENT2 = RGBColor(0xFF, 0x78, 0x47)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0xA8, 0xB3, 0xC5)
GREEN = RGBColor(0x34, 0xD3, 0x99)

FONT = "微软雅黑"


# =========================================================================
# PPT
# =========================================================================
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def new_slide():
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    r.fill.solid()
    r.fill.fore_color.rgb = BG
    r.line.fill.background()
    r.shadow.inherit = False
    return s


def box(slide, l, t, w, h, fill=None, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.adjustments[0] = 0.06
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(1)
    sh.shadow.inherit = False
    return sh


def txt(slide, l, t, w, h, runs, size=18, color=WHITE, bold=False,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, leading=1.15):
    """runs: str 或 [(text,{size,color,bold}), ...] 或多段 [[..],[..]]"""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    paras = runs if isinstance(runs, list) and runs and isinstance(runs[0], list) else [runs]
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = leading
        p.space_after = Pt(6)
        items = para if isinstance(para, list) else [(para, {})]
        for text, st in items:
            r = p.add_run()
            r.text = text
            f = r.font
            f.name = st.get("font", FONT)
            f.size = Pt(st.get("size", size))
            f.bold = st.get("bold", bold)
            f.color.rgb = st.get("color", color)
    return tb


def header(slide, kicker, title):
    box(slide, Inches(0.55), Inches(0.42), Inches(0.09), Inches(0.86), fill=ACCENT)
    txt(slide, Inches(0.85), Inches(0.38), Inches(11.9), Inches(0.34),
        kicker, size=13, color=ACCENT, bold=True)
    txt(slide, Inches(0.85), Inches(0.68), Inches(11.9), Inches(0.62),
        title, size=30, color=WHITE, bold=True)


def footer(slide, n):
    txt(slide, Inches(0.55), Inches(7.06), Inches(6), Inches(0.3),
        "token-burner · 文档消耗器", size=10, color=GRAY)
    txt(slide, Inches(12.3), Inches(7.06), Inches(0.55), Inches(0.3),
        str(n), size=10, color=GRAY, align=PP_ALIGN.RIGHT)


# ---- S1 封面 ----
s = new_slide()
box(s, Inches(0), Inches(0), SW, SH, fill=BG)
deco = box(s, Inches(9.2), Inches(-1.6), Inches(7), Inches(7), fill=PANEL)
fire = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.6), Inches(1.6),
                          Inches(3.4), Inches(3.4))
fire.fill.solid(); fire.fill.fore_color.rgb = ACCENT2
fire.line.fill.background(); fire.shadow.inherit = False
txt(s, Inches(10.95), Inches(2.6), Inches(2.7), Inches(1.4),
    [[("点", {"size": 48, "bold": True, "color": BG})]],
    align=PP_ALIGN.CENTER)
txt(s, Inches(0.9), Inches(2.1), Inches(8.6), Inches(0.4),
    "AI MULTI-AGENT PROJECT TEAM SYSTEM", size=14, color=ACCENT, bold=True)
txt(s, Inches(0.9), Inches(2.55), Inches(9.2), Inches(1.7),
    [[("token-burner", {"size": 54, "bold": True, "color": WHITE}),
      (" 文档消耗器", {"size": 40, "bold": True, "color": ACCENT})]])
txt(s, Inches(0.9), Inches(4.15), Inches(8.8), Inches(0.9),
    "让每个 token 都烧在刀刃上 —— 成本可控、质量可验证的多智能体软件开发系统",
    size=18, color=GRAY)
box(s, Inches(0.9), Inches(5.35), Inches(2.6), Inches(0.55), fill=ACCENT)
txt(s, Inches(0.9), Inches(5.44), Inches(2.6), Inches(0.4),
    "投资人汇报 · 2026.08", size=15, color=BG, bold=True,
    align=PP_ALIGN.CENTER)
txt(s, Inches(0.9), Inches(6.6), Inches(9), Inches(0.4),
    "规格 v0.3.1 全量实现 · 472 项自动化测试通过 · 桌面客户端已交付",
    size=14, color=GRAY)

# ---- S2 痛点 ----
s = new_slide()
header(s, "01 · 问题", "AI 编码工具的三大痛点")
pains = [
    ("成本失控", "Token 消耗不可预估、不可观测。讨论发散、修复死循环、"
     "单次输出倾泻——账单在无声中膨胀，企业级使用无法预算管理。", "💰"),
    ("能力天花板", "单模型自我校验不可靠：自己写的代码自己审，"
     "盲区一致、幻觉互相印证，缺乏独立视角的交叉验证。", "🎯"),
    ("质量无闸门", "生成代码直接交付：无接口契约、无静态验证、"
     "无执行反馈闭环，「看起来对」不等于「跑得对」。", "🚧"),
]
for i, (t1, t2, _) in enumerate(pains):
    x = Inches(0.55 + i * 4.18)
    box(s, x, Inches(1.95), Inches(3.9), Inches(3.6), fill=PANEL)
    box(s, x, Inches(1.95), Inches(3.9), Inches(0.12), fill=ACCENT2)
    txt(s, x + Inches(0.35), Inches(2.35), Inches(3.2), Inches(0.5),
        t1, size=22, bold=True, color=ACCENT)
    txt(s, x + Inches(0.35), Inches(3.0), Inches(3.25), Inches(2.4),
        t2, size=14.5, color=GRAY, leading=1.35)
box(s, Inches(0.55), Inches(5.9), Inches(12.23), Inches(0.85), fill=PANEL)
txt(s, Inches(0.9), Inches(6.08), Inches(11.6), Inches(0.55),
    [[("核心矛盾：", {"bold": True, "color": ACCENT, "size": 16}),
      ("多智能体能解决质量问题，但会放大成本问题——市场上缺少一个"
       "「既要多模型协作、又要成本铁闸」的编排层。", {"color": WHITE, "size": 16})]])
footer(s, 2)

# ---- S3 解决方案 ----
s = new_slide()
header(s, "02 · 解决方案", "一支被工程纪律约束的 AI 项目团队")
txt(s, Inches(0.55), Inches(1.62), Inches(12.2), Inches(0.5),
    "需求 → 评估路由 → 团队组建 → 方案讨论 → spec 收敛 → 模块化生成 → 双模验证 → 反馈修复",
    size=15, color=ACCENT, bold=True)
roles = [
    ("主 LLM · 架构师", "需求解析、难度评估、方案提出与汇总裁决、spec 生成"),
    ("开发副 LLM", "方案评审（≥8 分建议采纳）、模块代码生成、缺陷修复"),
    ("测试副 LLM", "方案评审、测试用例与测试文件生成、验证报告"),
]
for i, (t1, t2) in enumerate(roles):
    x = Inches(0.55 + i * 4.18)
    box(s, x, Inches(2.3), Inches(3.9), Inches(2.0), fill=PANEL)
    txt(s, x + Inches(0.32), Inches(2.6), Inches(3.3), Inches(0.45),
        t1, size=17, bold=True, color=WHITE)
    txt(s, x + Inches(0.32), Inches(3.12), Inches(3.3), Inches(1.1),
        t2, size=13.5, color=GRAY, leading=1.3)
feats = [
    "三模型强制互异——独立视角交叉验证，杜绝同源盲区",
    "模型无关接入（litellm）：OpenAI / 智谱 / 任意兼容端点自由组合",
    "简单任务节流：基础/研究类直出，不浪费团队编制",
]
for i, f in enumerate(feats):
    y = Inches(4.65 + i * 0.62)
    box(s, Inches(0.55), y, Inches(0.42), Inches(0.42), fill=ACCENT)
    txt(s, Inches(0.55), y + Inches(0.03), Inches(0.42), Inches(0.36),
        str(i + 1), size=16, bold=True, color=BG, align=PP_ALIGN.CENTER)
    txt(s, Inches(1.2), y + Inches(0.04), Inches(11.5), Inches(0.5),
        f, size=15.5, color=WHITE)
footer(s, 3)

# ---- S4 管线 ----
s = new_slide()
header(s, "03 · 工作管线", "三分类路由：把 token 花在值得的地方")
box(s, Inches(0.55), Inches(1.85), Inches(3.2), Inches(1.05), fill=ACCENT)
txt(s, Inches(0.55), Inches(2.12), Inches(3.2), Inches(0.5),
    "用户需求", size=19, bold=True, color=BG, align=PP_ALIGN.CENTER)
routes = [
    ("基础任务", "快速直答\n（零团队成本）", GREEN),
    ("研究·分析", "深度直出\n（单模型+预算约束）", ACCENT),
    ("编程任务", "强制组建三模型团队\n（完整工程管线）", ACCENT2),
]
for i, (t1, t2, c) in enumerate(routes):
    y = Inches(3.35 + i * 1.15)
    box(s, Inches(0.55), y, Inches(3.2), Inches(0.95), fill=PANEL)
    box(s, Inches(0.55), y, Inches(0.1), Inches(0.95), fill=c)
    txt(s, Inches(0.85), y + Inches(0.12), Inches(2.8), Inches(0.4),
        t1, size=16, bold=True, color=WHITE)
    txt(s, Inches(0.85), y + Inches(0.5), Inches(2.9), Inches(0.4),
        t2.replace("\n", " · "), size=12.5, color=GRAY)
stages = ["方案讨论\n(双评审+收敛)", "spec 确认\n(≤3轮)", "模块化生成\n(接口契约)",
          "双模验证\n(静态+执行)", "反馈修复\n(≤5轮)", "交付+看板\n(可审计)"]
for i, st in enumerate(stages):
    x = Inches(4.35 + i * 1.48)
    box(s, x, Inches(3.0), Inches(1.32), Inches(2.1), fill=PANEL)
    box(s, x, Inches(3.0), Inches(1.32), Inches(0.09), fill=ACCENT)
    lines = st.split("\n")
    txt(s, x + Inches(0.08), Inches(3.45), Inches(1.16), Inches(1.2),
        [[(lines[0], {"size": 13.5, "bold": True, "color": WHITE})],
         [(lines[1], {"size": 10.5, "color": GRAY})]],
        align=PP_ALIGN.CENTER)
    if i < len(stages) - 1:
        txt(s, x + Inches(1.3), Inches(3.85), Inches(0.25), Inches(0.4),
            "›", size=22, bold=True, color=ACCENT)
box(s, Inches(4.35), Inches(5.5), Inches(8.42), Inches(0.95), fill=PANEL)
txt(s, Inches(4.65), Inches(5.68), Inches(7.9), Inches(0.65),
    [[("全程铁闸：", {"bold": True, "color": ACCENT, "size": 14.5}),
      ("每一步都受六层成本护栏约束（下页）；中断可恢复，会话可审计。",
       {"color": WHITE, "size": 14.5})]])
footer(s, 4)

# ---- S5 六层护栏 ----
s = new_slide()
header(s, "04 · 核心壁垒", "六层成本护栏：把「无限烧钱」变成「预算工程」")
rows = [
    ("层", "护栏", "默认参数", "达标行为"),
    ("0", "单任务总预算（总闸）", "200k tokens · 自动模式 ×2.5",
     "≥90% 省token模式；超预算中止并落盘止损"),
    ("1", "讨论轮数上限", "3 轮", "主 LLM 直接产出收敛 spec"),
    ("2", "单轮输出上限", "8k tokens", "截断按未完成处理，分块续写"),
    ("3", "讨论循环检测", "Jaccard 0.9 / embedding 0.85 / 3 次",
     "冻结副 LLM 发言权，主 LLM 收权裁决"),
    ("4", "修复循环上限", "5 轮", "输出「已知问题与降级方案」"),
    ("5", "spec 确认收敛", "3 次", "合并意见输出最终 spec"),
]
tbl = s.shapes.add_table(7, 4, Inches(0.55), Inches(1.85),
                         Inches(12.23), Inches(4.3)).table
tbl.columns[0].width = Inches(0.7)
tbl.columns[1].width = Inches(3.3)
tbl.columns[2].width = Inches(3.9)
tbl.columns[3].width = Inches(4.33)
for ri, row in enumerate(rows):
    for ci, val in enumerate(row):
        cell = tbl.cell(ri, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT if ri == 0 else (
            PANEL if ri % 2 else RGBColor(0x16, 0x20, 0x33))
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Inches(0.14)
        cell.margin_top = cell.margin_bottom = Inches(0.02)
        p = cell.text_frame.paragraphs[0]
        r = p.add_run(); r.text = val
        r.font.name = FONT
        r.font.size = Pt(13 if ri else 14)
        r.font.bold = ri == 0 or ci == 0
        r.font.color.rgb = BG if ri == 0 else (ACCENT if ci == 0 else WHITE)
txt(s, Inches(0.55), Inches(6.35), Inches(12.2), Inches(0.6),
    [[("差异化本质：", {"bold": True, "color": ACCENT, "size": 15}),
      ("护栏全部由程序确定性执行（AST/相似度计算），不依赖 LLM 自律——"
       "「让 LLM 自己省 token」是不成立的，规则必须由代码铸成。",
       {"color": GRAY, "size": 15})]])
footer(s, 5)

# ---- S6 质量保障 ----
s = new_slide()
header(s, "05 · 质量保障", "验证不是提示词，是机制")
cards = [
    ("双模执行验证",
     ["安全审阅模式：AST 静态检查 + LLM 逻辑审查 + 用户反馈",
      "自动验证模式：危险预扫描（系统命令/网络/动态执行）",
      "真实执行 + 30s 超时熔断，结构化结果注入修复"]),
    ("AST 接口门禁",
     ["接口地图为唯一合法基线，引用抽取差异判定",
      "违规即阻断进入修复循环",
      "零执行、零 LLM 调用——唯一不耗 token 的纪律执行器"]),
    ("循环检测双保险",
     ["第一道：Jaccard 文本相似度（零成本拦截字面重复）",
      "第二道：embedding 语义相似度（增量缓存控成本）",
      "重复 ≥3 次自动冻结，主 LLM 收权裁决"]),
    ("韧性工程",
     ["LLM 超时 120s + 指数退避重试（瞬态错误最多 3 次）",
      "中断恢复：状态落盘，重启后断点续跑",
      "提示词注入防护：不可信文本边界标记 + 截断"]),
]
for i, (t1, items) in enumerate(cards):
    x = Inches(0.55 + (i % 2) * 6.27)
    y = Inches(1.9 + (i // 2) * 2.42)
    box(s, x, y, Inches(5.95), Inches(2.2), fill=PANEL)
    txt(s, x + Inches(0.3), y + Inches(0.18), Inches(5.4), Inches(0.42),
        t1, size=17, bold=True, color=ACCENT)
    paras = [[("▪ ", {"color": ACCENT2, "size": 13}),
              (it, {"color": GRAY, "size": 13})] for it in items]
    txt(s, x + Inches(0.3), y + Inches(0.66), Inches(5.4), Inches(1.45),
        paras, leading=1.25)
footer(s, 6)

# ---- S7 工程完成度 ----
s = new_slide()
header(s, "06 · 工程完成度", "不是 PPT 工程，是可交付的软件")
stats = [("规格 v0.3.1", "全量实现"), ("472", "自动化测试通过"),
         ("3 种", "交付形态"), ("0", "已知阻断缺陷")]
for i, (v, k) in enumerate(stats):
    x = Inches(0.55 + i * 3.12)
    box(s, x, Inches(1.9), Inches(2.9), Inches(1.55), fill=PANEL)
    txt(s, x, Inches(2.1), Inches(2.9), Inches(0.7),
        v, size=34, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    txt(s, x, Inches(2.85), Inches(2.9), Inches(0.4),
        k, size=14, color=GRAY, align=PP_ALIGN.CENTER)
done = [
    "完整管线：路由/讨论/spec/模块化生成/双模验证/反馈闭环",
    "六层护栏 + Token 成本看板（任务级消耗可审计）",
    "模块化治理：接口隔离、独立规格 md、共享库变更整包回归",
    "JSON 解析容错：围栏剥离、截断分块续写、LLM 辅助修复",
    "安全与信任边界文档化：注入防护、执行隔离、权限声明",
]
txt(s, Inches(0.55), Inches(3.85), Inches(7.5), Inches(0.4),
    "核心能力清单", size=16, bold=True, color=WHITE)
for i, d in enumerate(done):
    y = Inches(4.35 + i * 0.5)
    txt(s, Inches(0.55), y, Inches(0.35), Inches(0.4), "✓",
        size=15, bold=True, color=GREEN)
    txt(s, Inches(0.95), y, Inches(7.1), Inches(0.45), d, size=13.5, color=GRAY)
box(s, Inches(8.85), Inches(3.85), Inches(3.93), Inches(2.9), fill=PANEL)
txt(s, Inches(9.15), Inches(4.05), Inches(3.4), Inches(0.4),
    "质量方法论", size=16, bold=True, color=ACCENT)
txt(s, Inches(9.15), Inches(4.5), Inches(3.4), Inches(2.1),
    "SDD 规格驱动开发 + TDD 测试驱动开发。\n\n规格先行、测试先行，"
    "每项缺陷修复均回归全量测试套件，472 项测试构成持续验证的安全网。",
    size=13, color=GRAY, leading=1.35)
footer(s, 7)

# ---- S8 产品形态 ----
s = new_slide()
header(s, "07 · 产品形态", "双击即用的桌面客户端 + 开发者生态")
forms = [
    ("桌面客户端", "独立 exe（75MB），双击即用。需求输入 → 点火执行 → "
     "工作可视化（评估指标/时间线/模块状态）→ API 配置管理。"
     "单进程零网络架构，兼容企业安全环境。"),
    ("Web / API 服务", "FastAPI 契约化接口：评估路由、任务执行、中断恢复、"
     "反馈闭环、成本看板。支持浏览器访问与二次集成。"),
    ("CLI 命令行", "完整管线命令行形态，面向开发者的脚本化、"
     "CI 集成场景，与 GUI 共用同一内核。"),
]
for i, (t1, t2) in enumerate(forms):
    y = Inches(1.9 + i * 1.55)
    box(s, Inches(0.55), y, Inches(12.23), Inches(1.35), fill=PANEL)
    box(s, Inches(0.55), y, Inches(0.12), Inches(1.35), fill=ACCENT)
    txt(s, Inches(0.95), y + Inches(0.16), Inches(2.6), Inches(0.5),
        t1, size=18, bold=True, color=WHITE)
    txt(s, Inches(0.95), y + Inches(0.62), Inches(11.3), Inches(0.7),
        t2, size=13.5, color=GRAY, leading=1.3)
txt(s, Inches(0.55), Inches(6.6), Inches(12.2), Inches(0.5),
    [[("一处内核，三种入口：", {"bold": True, "color": ACCENT, "size": 14.5}),
      ("桌面端覆盖小白用户，API/CLI 覆盖开发者与企业集成——同一套多智能体内核。",
       {"color": GRAY, "size": 14.5})]])
footer(s, 8)

# ---- S9 市场 ----
s = new_slide()
header(s, "08 · 市场机会", "AI 编码是生成式 AI 落地最快的赛道")
txt(s, Inches(0.55), Inches(1.85), Inches(12.2), Inches(1.0),
    "从「单点补全」到「团队交付」：开发者工具的下一个范式转移已经开始。"
    "GitHub Copilot、Cursor 验证了单模型辅助的巨大市场；而多智能体编排 + "
    "成本治理层，是尚未被占领的下一层价值。",
    size=15.5, color=WHITE, leading=1.4)
segs = [
    ("企业 AI 编码预算管理", "企业引入 AI 编程的最大阻力从「能力」转向「成本失控」。"
     "预算铁闸 + 成本看板直接命中采购决策者的核心关切。"),
    ("多模型异构编排", "模型市场碎片化（OpenAI/Anthropic/国产模型各有所长），"
     "「三模型互异交叉验证」天然成为模型中立的编排层。"),
    ("交付质量合规", "代码生成进入企业生产链路需要验证闭环与审计留痕——"
     "本项目接口门禁 + 会话记录 + 成本审计即为此设计。"),
]
for i, (t1, t2) in enumerate(segs):
    x = Inches(0.55 + i * 4.18)
    box(s, x, Inches(3.15), Inches(3.9), Inches(2.5), fill=PANEL)
    txt(s, x + Inches(0.3), Inches(3.42), Inches(3.3), Inches(0.75),
        t1, size=16.5, bold=True, color=ACCENT, leading=1.2)
    txt(s, x + Inches(0.3), Inches(4.2), Inches(3.3), Inches(1.4),
        t2, size=13, color=GRAY, leading=1.3)
txt(s, Inches(0.55), Inches(6.1), Inches(12.2), Inches(0.7),
    [[("定位一句话：", {"bold": True, "color": ACCENT2, "size": 15}),
      ("不做又一个编码助手，做 AI 软件团队的「操作系统 + 财务总监」。",
       {"color": WHITE, "size": 15, "bold": True})]])
footer(s, 9)

# ---- S10 竞争壁垒 ----
s = new_slide()
header(s, "09 · 竞争壁垒", "机制壁垒，而非提示词壁垒")
moats = [
    ("护栏引擎", "六层确定性护栏 + 超预算熔断落盘。竞品的「成本提示」"
     "是建议，本项目的是闸门。"),
    ("AST 门禁", "接口地图 + 零成本确定性校验，把「开发纪律」变成"
     "「硬门禁」——提示词无法被绕过。"),
    ("模型异构互验", "强制三模型互异，交叉盲区检测。单模型产品结构上"
     "无法复制。"),
    ("全链路可审计", "会话记录、成本报告、修复历史、中断快照——"
     "每一步留痕，企业级合规友好。"),
]
for i, (t1, t2) in enumerate(moats):
    x = Inches(0.55 + (i % 2) * 6.27)
    y = Inches(1.95 + (i // 2) * 2.15)
    box(s, x, y, Inches(5.95), Inches(1.95), fill=PANEL)
    box(s, x + Inches(0.3), y + Inches(0.3), Inches(0.55), Inches(0.55),
        fill=ACCENT)
    txt(s, x + Inches(0.3), y + Inches(0.38), Inches(0.55), Inches(0.45),
        str(i + 1), size=20, bold=True, color=BG, align=PP_ALIGN.CENTER)
    txt(s, x + Inches(1.05), y + Inches(0.32), Inches(4.6), Inches(0.5),
        t1, size=18, bold=True, color=WHITE)
    txt(s, x + Inches(1.05), y + Inches(0.88), Inches(4.7), Inches(0.95),
        t2, size=13.5, color=GRAY, leading=1.3)
txt(s, Inches(0.55), Inches(6.45), Inches(12.2), Inches(0.5),
    "以上能力均有规格文档锚定 + 自动化测试保护，构成可验证的工程壁垒。",
    size=14, color=GRAY)
footer(s, 10)

# ---- S11 路线图 ----
s = new_slide()
header(s, "10 · 路线图", "MVP 已交付，迭代路径清晰")
phases = [
    ("MVP · 已交付", "2026 Q3", GREEN,
     ["完整协作闭环 + 双模验证", "六层护栏 + 成本看板", "桌面客户端 / API / CLI",
      "472 项测试、规格全实现"]),
    ("Alpha v0.4", "2026 Q4", ACCENT,
     ["沙箱隔离执行（本地执行器已就绪）", "外部工具 / MCP 调用开放",
      "模式状态智能推荐"]),
    ("Beta v0.5", "2027 H1", ACCENT2,
     ["Researcher Agent（联网调研）", "实时监控面板", "Agent 对话流图可视化"]),
    ("Release v1.0", "2027 H2", WHITE,
     ["本地知识库检索", "副 LLM 能力自动匹配", "多人协作与权限管理"]),
]
for i, (t1, t2, c, items) in enumerate(phases):
    x = Inches(0.55 + i * 3.12)
    box(s, x, Inches(1.95), Inches(2.9), Inches(4.3), fill=PANEL)
    box(s, x, Inches(1.95), Inches(2.9), Inches(0.14), fill=c)
    txt(s, x + Inches(0.25), Inches(2.28), Inches(2.4), Inches(0.45),
        t1, size=17, bold=True, color=c)
    txt(s, x + Inches(0.25), Inches(2.72), Inches(2.4), Inches(0.35),
        t2, size=12.5, color=GRAY)
    paras = [[("▪ ", {"color": c, "size": 12}),
              (it, {"color": GRAY, "size": 12})] for it in items]
    txt(s, x + Inches(0.25), Inches(3.25), Inches(2.45), Inches(2.8),
        paras, leading=1.3)
footer(s, 11)

# ---- S12 尾页 ----
s = new_slide()
box(s, Inches(0), Inches(2.3), SW, Inches(2.9), fill=PANEL)
txt(s, Inches(0), Inches(2.75), SW, Inches(0.9),
    "Set your tokens on fire — 有价值地燃烧", size=34, bold=True,
    color=WHITE, align=PP_ALIGN.CENTER)
txt(s, Inches(0), Inches(3.85), SW, Inches(0.6),
    "成本可控 · 质量可验证 · 全程可审计的 AI 多智能体软件工厂",
    size=18, color=ACCENT, align=PP_ALIGN.CENTER)
txt(s, Inches(0), Inches(5.8), SW, Inches(0.5),
    "谢谢 · 期待与您深入交流", size=16, color=GRAY, align=PP_ALIGN.CENTER)

pptx_path = OUT / "token-burner_投资人汇报.pptx"
prs.save(pptx_path)
print("PPT saved:", pptx_path)


# =========================================================================
# 可行性研究报告 (docx)
# =========================================================================
doc = Document()
style = doc.styles["Normal"]
style.font.name = "Calibri"
style.font.size = Pt(11)
style.element.rPr.rFonts.set(qn("w:eastAsia"), FONT)

NAVY = RGBColor(0x10, 0x18, 0x28)


def h(level, text_):
    p = doc.add_heading(text_, level)
    for r in p.runs:
        r.font.name = FONT
        r._element.rPr.rFonts.set(qn("w:eastAsia"), FONT)
    return p


def para(text_, bold=False, size=11, color=None):
    p = doc.add_paragraph()
    r = p.add_run(text_)
    r.font.name = FONT
    r._element.rPr.rFonts.set(qn("w:eastAsia"), FONT)
    r.font.size = Pt(size)
    r.font.bold = bold
    return p


def bullet(text_):
    p = doc.add_paragraph(style="List Bullet")
    r = p.add_run(text_)
    r.font.name = FONT
    r._element.rPr.rFonts.set(qn("w:eastAsia"), FONT)
    r.font.size = Pt(11)
    return p


def table(rows, widths=None):
    t = doc.add_table(rows=len(rows), cols=len(rows[0]))
    t.style = "Table Grid"
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = t.cell(ri, ci)
            cell.text = ""
            p = cell.paragraphs[0]
            r = p.add_run(val)
            r.font.name = FONT
            r._element.rPr.rFonts.set(qn("w:eastAsia"), FONT)
            r.font.size = Pt(9.5)
            r.font.bold = ri == 0
    if widths:
        for ci, w in enumerate(widths):
            for ri in range(len(rows)):
                t.cell(ri, ci).width = Inches(w)
    doc.add_paragraph()
    return t


# 封面
for _ in range(5):
    doc.add_paragraph()
p = para("token-burner · 文档消耗器", bold=True, size=28)
p.alignment = WD_ALIGN_PARAGRAPH.CENTER
p = para("AI 多智能体项目团队系统", size=20)
p.alignment = WD_ALIGN_PARAGRAPH.CENTER
p = para("可行性研究报告", bold=True, size=24)
p.alignment = WD_ALIGN_PARAGRAPH.CENTER
doc.add_paragraph()
p = para("（投资人版 · 2026 年 8 月）", size=13)
p.alignment = WD_ALIGN_PARAGRAPH.CENTER
doc.add_page_break()

# 1 执行摘要
h(1, "一、执行摘要")
para("token-burner（文档消耗器）是一套 AI 多智能体项目团队系统：以「主 LLM 架构师 + "
     "开发副 LLM + 测试副 LLM」的固定团队编制，自动完成「需求 → 方案讨论 → spec 生成 → "
     "模块化代码生成 → 双模验证 → 反馈修复」的软件开发闭环。产品以规格文档 v0.3.1 为唯一"
     "需求基线，采用 SDD（规格驱动）+ TDD（测试驱动）方法论开发，当前已完成 MVP 全量功能，"
     "472 项自动化测试通过，桌面客户端（独立 exe）、Web API、CLI 三种形态均可交付。")
para("与市面单模型编码助手相比，本产品的核心差异化在于：①六层确定性成本护栏（预算总闸、"
     "讨论轮数、单轮输出、循环检测、修复上限、spec 收敛），将 token 消耗从「不可预估」"
     "变为「预算工程」；②三模型强制互异的交叉验证与 AST 接口门禁，将「提示词纪律」升级为"
     "「机制闸门」；③全链路留痕（会话记录、成本报告、修复历史、中断快照），天然满足企业级"
     "审计与合规诉求。")
para("本报告从技术、产品、经济、风险四个维度论证：项目已具备从 MVP 走向商业化的可行性。", bold=True)

# 2 背景
h(1, "二、项目背景与市场分析")
h(2, "2.1 行业痛点")
bullet("成本失控：多智能体系统通过多轮对话与多模型协作提升质量，但同时放大 token 消耗；"
       "讨论发散、修复死循环、单次输出倾泻等场景使账单难以预算，成为企业规模采用的现实阻力。")
bullet("能力天花板：单模型「自己写、自己审」缺乏独立视角，盲区一致、幻觉互相印证；"
       "模型市场碎片化背景下，单一模型难以在架构、实现、测试三视角上同时占优。")
bullet("质量无闸门：主流工具生成的代码缺乏接口契约、静态验证与执行反馈闭环，"
       "「看起来对」不等于「跑得对」，难以进入企业生产链路。")
h(2, "2.2 市场定位")
para("AI 编码是生成式 AI 商业化落地最快、付费意愿最强的赛道之一。GitHub Copilot、Cursor "
     "等产品验证了「单模型辅助编程」的市场规模；随着企业关注点从「能不能用」转向「成本是否"
     "可控、质量是否可验证、过程是否可审计」，多智能体编排与成本治理层构成尚未被占领的下一层"
     "价值。本产品的定位是：")
para("「不做又一个编码助手，做 AI 软件团队的操作系统 + 财务总监。」", bold=True)
h(2, "2.3 目标客群")
bullet("中小软件团队与外包公司：需要降低开发成本、但无法承担 AI 账单失控风险。")
bullet("企业 IT 部门：需要可审计、可预算管理的 AI 编码工具以满足合规要求。")
bullet("AI 工具集成商与开发者：通过 API/CLI 将多智能体内核嵌入自有工作流。")

# 3 技术可行性
h(1, "三、技术可行性")
h(2, "3.1 系统架构")
para("系统采用分层架构：交互层（桌面客户端 / Web API / CLI，共用同一契约）→ 编排层"
     "（任务路由、团队组建、方案讨论、管道管理）→ 智能体层（开发循环、模块构建、"
     "反馈修复）→ 基础设施层（模型客户端、预算闸门、相似度检测、静态/接口校验、"
     "执行器、成本看板）。模型接入基于 litellm 统一抽象，OpenAI、智谱等任意兼容端点"
     "可自由组合，避免单一厂商绑定。")
h(2, "3.2 六层成本护栏（核心技术壁垒）")
table([
    ("层", "护栏", "默认参数", "达标行为"),
    ("0", "单任务总预算（总闸）", "200k tokens（自动模式 ×2.5）",
     "≥90% 进入省 token 模式；超预算中止并落盘已完成部分与止损清单"),
    ("1", "讨论轮数上限", "3 轮", "主 LLM 直接产出收敛 spec"),
    ("2", "单轮输出上限", "8k tokens", "截断按未完成处理，分块续写"),
    ("3", "讨论循环检测", "Jaccard 0.9 / embedding 0.85 / 重复 3 次",
     "冻结副 LLM 发言权，主 LLM 收权裁决"),
    ("4", "修复循环上限", "5 轮", "输出「已知问题与降级方案」交用户决策"),
    ("5", "spec 确认收敛", "3 次", "合并意见输出最终 spec"),
], widths=[0.4, 1.7, 2.0, 2.9])
para("关键设计原则：所有护栏由程序确定性执行（AST 解析、相似度计算、预算累加），"
     "不依赖 LLM 自律。这一原则使成本上界在任务开始前即可预估——这是企业采购决策的"
     "核心关切，也是单模型产品结构上无法复制的差异。")
h(2, "3.3 质量保障机制")
bullet("双模执行验证：安全审阅模式（AST 静态检查 + LLM 逻辑审查 + 用户手动反馈）与"
       "自动验证模式（危险预扫描 + 真实子进程执行 + 30 秒超时熔断）双轨并行，执行器抽象"
       "支持渐进式信任。")
bullet("AST 接口门禁：以接口地图为唯一合法基线做引用抽取与差异判定，违规即阻断——"
       "零执行、零 LLM 调用，是全系统唯一不消耗 token 的纪律执行器。")
bullet("循环检测双保险：Jaccard 文本相似度首道拦截（零成本）+ embedding 语义相似度"
       "第二道拦截（增量缓存），杜绝讨论空转烧钱。")
bullet("韧性工程：LLM 调用超时（默认 120s）与指数退避重试（最多 3 次）；崩溃中断后"
       "状态落盘，支持断点续跑；不可信文本注入提示词前强制边界标记与截断。")
h(2, "3.4 已验证的工程成果")
table([
    ("维度", "成果", "证据"),
    ("功能完成度", "规格 v0.3.1 全量实现（路由/讨论/生成/验证/反馈/看板/恢复）",
     "规格条款逐项比对审计通过"),
    ("测试覆盖", "472 项自动化测试通过、1 项跳过、0 失败",
     "pytest 全量回归"),
    ("交付形态", "桌面 exe（75MB，双击即用）/ Web API / CLI", "三形态共用同一内核"),
    ("规格外加固", "8 项深度审计问题全部修复（可运行性/超时重试/中断恢复/注入防护等）",
     "审计-修复记录闭环"),
], widths=[1.1, 3.2, 2.7])
para("结论：技术方案全部落地且经自动化测试验证，无未解决的技术阻塞项。", bold=True)

# 4 产品可行性
h(1, "四、产品可行性")
h(2, "4.1 产品形态与用户体验")
para("桌面客户端面向非技术用户：左侧需求输入与「点火」执行，右侧 LLM 工作可视化"
     "（评估指标、执行时间线、模块状态）与 API 接口配置管理；采用单进程零网络架构，"
     "无需安装 Python 环境，双击即用，且对企业安全软件环境兼容。Web/API 与 CLI 形态"
     "服务开发者与企业集成场景。")
h(2, "4.2 开发方法论带来的交付确定性")
para("项目全程采用 SDD + TDD：任何功能先有规格条款、后有实现与对应测试。这使得"
     "「产品是否做完」成为一个可机检的问题而非主观判断——472 项测试即交付质量的"
     "持续证明，也为后续迭代提供回归安全网。")

# 5 经济可行性
h(1, "五、经济可行性")
h(2, "5.1 成本结构")
para("产品边际成本主要是 LLM API 的 token 费用。以默认预算 200k tokens/任务计"
     "（自动验证模式 ×2.5 = 500k），按主流模型公开定价区间粗略估算（输入约 $0.5–$3/"
     "百万 token、输出约 $1.5–$15/百万 token，具体以各厂商实时费率为准），单任务直接"
     "成本约为零点几美元至数美元量级，且在任务启动前即可预估上界——这本身就是产品"
     "核心能力对自身成本的复用。")
h(2, "5.2 商业化路径")
bullet("订阅制（SaaS）：按任务预算档位订阅，护栏与看板为付费核心。")
bullet("私有化部署：面向数据合规敏感的企业，按席位/年授权。")
bullet("API 计量：面向集成商按调用量计费，多智能体内核作为 PnP 组件。")
para("注：以上为方向性框架，具体定价需在 Beta 阶段以真实用量数据校准。")

# 6 风险
h(1, "六、风险分析与对策")
table([
    ("风险", "等级", "对策"),
    ("模型能力迭代导致编排层价值被稀释", "中",
     "壁垒锚定在「确定性护栏 + 审计合规」而非模型能力；模型越强，"
     "企业对成本与质量治理的需求反而越大"),
    ("LLM 费率波动影响毛利", "中",
     "模型无关架构支持费率优化切换；预算工程化能力同时服务自身成本控制"),
    ("安全软件误报（未签名 exe）", "低",
     "已通过单进程零网络架构规避；正式发布将配合代码签名证书"),
    ("竞品跟进多智能体编排", "中",
     "六层护栏、AST 门禁、可审计链路均有规格与测试双保护，"
     "工程壁垒复制成本高"),
    ("提示词注入等安全对抗", "低",
     "已实现不可信文本边界隔离与截断；信任边界已在文档中显式声明"),
], widths=[1.9, 0.6, 4.5])

# 7 结论
h(1, "七、结论与建议")
para("token-burner 已完成规格 v0.3.1 定义的全部 MVP 能力，并以 472 项自动化测试、"
     "三种交付形态与多轮深度审计修复证明其工程成熟度。技术上无阻塞项，产品上具备"
     "清晰差异化（成本护栏 + 机制化质量门禁 + 全链路审计），经济上边际成本可控且"
     "可预估。")
para("建议：本项目具备进入 Alpha/Beta 迭代并启动早期商业化验证的条件。", bold=True)

docx_path = OUT / "token-burner_可行性研究报告.docx"
doc.save(docx_path)
print("DOCX saved:", docx_path)
